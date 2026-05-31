import os
import pptx
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

def create_presentation():
    prs = pptx.Presentation()
    # Set slide dimensions to 16:9 widescreen
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    
    # Custom Theme Colors
    DARK_BG = RGBColor(10, 25, 47)      # Dark Navy Blue (Premium Background)
    TEXT_LIGHT = RGBColor(240, 245, 255) # Off-white for text
    ACCENT_CYAN = RGBColor(0, 230, 118)  # Neon Mint Green / Cyan accent
    ACCENT_BLUE = RGBColor(27, 38, 59)   # Medium slate blue for containers
    TEXT_GRAY = RGBColor(180, 190, 210)  # Muted grey-blue for secondary text
    
    # ----------------------------------------------------
    # Helper to style background
    # ----------------------------------------------------
    def apply_dark_background(slide):
        # Draw a full slide rectangle to apply background color
        background = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
        )
        background.fill.solid()
        background.fill.fore_color.rgb = DARK_BG
        background.line.fill.background() # No line
        
        # Lower shape send to back
        slide.shapes._spTree.remove(background._element)
        slide.shapes._spTree.insert(2, background._element)
        return background

    # Helper to add standard header and footer
    def add_slide_header(slide, title_text):
        apply_dark_background(slide)
        
        # Header Box
        header_box = slide.shapes.add_textbox(Inches(0.75), Inches(0.5), Inches(11.83), Inches(0.8))
        tf = header_box.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
        
        p = tf.paragraphs[0]
        p.text = title_text
        p.font.name = "Trebuchet MS"
        p.font.size = Pt(36)
        p.font.bold = True
        p.font.color.rgb = ACCENT_CYAN
        
        # Add subtle accent line
        line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0.75), Inches(1.35), Inches(3.0), Inches(0.04)
        )
        line.fill.solid()
        line.fill.fore_color.rgb = ACCENT_CYAN
        line.line.fill.background()
        
        # Footer text
        footer_box = slide.shapes.add_textbox(Inches(0.75), Inches(7.0), Inches(11.83), Inches(0.3))
        ftf = footer_box.text_frame
        ftf.word_wrap = True
        fp = ftf.paragraphs[0]
        fp.text = "AI-Powered Career Path Recommendation (CareerCraft) | BE-CSE"
        fp.font.name = "Arial"
        fp.font.size = Pt(10)
        fp.font.color.rgb = TEXT_GRAY

    # ----------------------------------------------------
    # SLIDE 1: Title Slide (Dark Theme, Stunning Design)
    # ----------------------------------------------------
    blank_slide_layout = prs.slide_layouts[6]
    slide1 = prs.slides.add_slide(blank_slide_layout)
    apply_dark_background(slide1)
    
    # Decorative glassmorphism card
    card = slide1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.75), Inches(0.75), Inches(11.83), Inches(6.0))
    card.fill.solid()
    card.fill.fore_color.rgb = ACCENT_BLUE
    card.line.color.rgb = ACCENT_CYAN
    card.line.width = Pt(1.5)
    
    # Title & Subtitle box inside the card
    title_box = slide1.shapes.add_textbox(Inches(1.2), Inches(1.2), Inches(11.0), Inches(2.2))
    tf = title_box.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "AI-POWERED CAREER PATH RECOMMENDATION"
    p.font.name = "Trebuchet MS"
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = ACCENT_CYAN
    p.alignment = PP_ALIGN.LEFT
    
    p2 = tf.add_paragraph()
    p2.text = "CareerCraft: Your Ultimate Career Compass & B2B Recruitment Platform"
    p2.font.name = "Arial"
    p2.font.size = Pt(18)
    p2.font.color.rgb = TEXT_LIGHT
    p2.alignment = PP_ALIGN.LEFT
    p2.space_before = Pt(10)
    
    # Student Details Box
    student_box = slide1.shapes.add_textbox(Inches(1.2), Inches(3.6), Inches(5.0), Inches(2.5))
    tf_student = student_box.text_frame
    tf_student.word_wrap = True
    
    ps1 = tf_student.paragraphs[0]
    ps1.text = "Submitted By:"
    ps1.font.name = "Arial"
    ps1.font.size = Pt(14)
    ps1.font.bold = True
    ps1.font.color.rgb = ACCENT_CYAN
    
    ps2 = tf_student.add_paragraph()
    ps2.text = "HARISHGANTH.L\nReg No: 720823103057\nFinal Year, BE-CSE"
    ps2.font.name = "Arial"
    ps2.font.size = Pt(16)
    ps2.font.bold = True
    ps2.font.color.rgb = TEXT_LIGHT
    ps2.space_before = Pt(5)
    
    # College and Guide Box
    college_box = slide1.shapes.add_textbox(Inches(6.5), Inches(3.6), Inches(5.5), Inches(2.5))
    tf_college = college_box.text_frame
    tf_college.word_wrap = True
    
    pc1 = tf_college.paragraphs[0]
    pc1.text = "Project Guide:"
    pc1.font.name = "Arial"
    pc1.font.size = Pt(14)
    pc1.font.bold = True
    pc1.font.color.rgb = ACCENT_CYAN
    
    pc2 = tf_college.add_paragraph()
    pc2.text = "Mrs. VIDHYA.V\nAssistant Professor, Dept of CSE"
    pc2.font.name = "Arial"
    pc2.font.size = Pt(16)
    pc2.font.color.rgb = TEXT_LIGHT
    pc2.space_before = Pt(5)
    
    pc3 = tf_college.add_paragraph()
    pc3.text = "\nHINDUSTAN INSTITUTE OF TECHNOLOGY"
    pc3.font.name = "Arial"
    pc3.font.size = Pt(14)
    pc3.font.bold = True
    pc3.font.color.rgb = TEXT_LIGHT
    
    # ----------------------------------------------------
    # SLIDE 2: Introduction & Problem Statement
    # ----------------------------------------------------
    slide2 = prs.slides.add_slide(blank_slide_layout)
    add_slide_header(slide2, "Introduction & Problem Statement")
    
    # Content left (Introduction)
    left_box = slide2.shapes.add_textbox(Inches(0.75), Inches(1.8), Inches(5.6), Inches(4.5))
    tf_left = left_box.text_frame
    tf_left.word_wrap = True
    
    p = tf_left.paragraphs[0]
    p.text = "Introduction"
    p.font.name = "Arial"
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = ACCENT_CYAN
    
    bullets = [
        "CareerCraft is an AI-powered 'Super App' designed for career exploration and B2B recruitment.",
        "Integrates dynamic machine learning guidance with resume analysis, interview coaching, and project suggestions.",
        "Acts as a comprehensive roadmap for graduating students to achieve industry readiness."
    ]
    for b in bullets:
        p = tf_left.add_paragraph()
        p.text = "• " + b
        p.font.name = "Arial"
        p.font.size = Pt(15)
        p.font.color.rgb = TEXT_LIGHT
        p.space_before = Pt(12)
        
    # Content right (Problem Statement)
    right_box = slide2.shapes.add_textbox(Inches(6.9), Inches(1.8), Inches(5.6), Inches(4.5))
    tf_right = right_box.text_frame
    tf_right.word_wrap = True
    
    p = tf_right.paragraphs[0]
    p.text = "Problem Statement"
    p.font.name = "Arial"
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = ACCENT_CYAN
    
    problems = [
        "Students face massive confusion choosing structured career roadmaps (Full Stack, Data Science, Cyber).",
        "Standard resumes lack ATS-friendliness, leading to high rejection rates.",
        "Lack of realistic, feedback-oriented technical mock interview practice.",
        "Recruiters struggle to efficiently screen, find, and match candidates who completed verified learning milestones."
    ]
    for pr in problems:
        p = tf_right.add_paragraph()
        p.text = "• " + pr
        p.font.name = "Arial"
        p.font.size = Pt(15)
        p.font.color.rgb = TEXT_LIGHT
        p.space_before = Pt(12)

    # ----------------------------------------------------
    # SLIDE 3: Proposed System Architecture
    # ----------------------------------------------------
    slide3 = prs.slides.add_slide(blank_slide_layout)
    add_slide_header(slide3, "System Architecture")
    
    # Large container
    container = slide3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.75), Inches(1.8), Inches(11.83), Inches(4.8))
    container.fill.solid()
    container.fill.fore_color.rgb = ACCENT_BLUE
    container.line.color.rgb = ACCENT_CYAN
    
    # Left column (Client Side UI)
    box1 = slide3.shapes.add_textbox(Inches(1.0), Inches(2.2), Inches(3.2), Inches(4.0))
    tf1 = box1.text_frame
    tf1.word_wrap = True
    p1 = tf1.paragraphs[0]
    p1.text = "1. Frontend (React 19)"
    p1.font.name = "Arial"
    p1.font.size = Pt(18)
    p1.font.bold = True
    p1.font.color.rgb = ACCENT_CYAN
    
    f_bullets = [
        "Vite-powered Single Page App",
        "Neon Glassmorphism Clean UI",
        "React Router 7 Navigation",
        "Custom vanilla CSS layouts",
        "Responsive Mobile-friendly"
    ]
    for b in f_bullets:
        p = tf1.add_paragraph()
        p.text = "⚡ " + b
        p.font.name = "Arial"
        p.font.size = Pt(14)
        p.font.color.rgb = TEXT_LIGHT
        p.space_before = Pt(8)
        
    # Middle column (Backend API)
    box2 = slide3.shapes.add_textbox(Inches(4.8), Inches(2.2), Inches(3.5), Inches(4.0))
    tf2 = box2.text_frame
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    p2.text = "2. Backend (Node & Express)"
    p2.font.name = "Arial"
    p2.font.size = Pt(18)
    p2.font.bold = True
    p2.font.color.rgb = ACCENT_CYAN
    
    b_bullets = [
        "RESTful API endpoints",
        "MongoDB schemas & models",
        "JWT Session Authentication",
        "Secured environment keys",
        "Error-handling controllers"
    ]
    for b in b_bullets:
        p = tf2.add_paragraph()
        p.text = "⚙️ " + b
        p.font.name = "Arial"
        p.font.size = Pt(14)
        p.font.color.rgb = TEXT_LIGHT
        p.space_before = Pt(8)
        
    # Right column (AI Engine & Db)
    box3 = slide3.shapes.add_textbox(Inches(9.0), Inches(2.2), Inches(3.2), Inches(4.0))
    tf3 = box3.text_frame
    tf3.word_wrap = True
    p3 = tf3.paragraphs[0]
    p3.text = "3. AI Engine & Database"
    p3.font.name = "Arial"
    p3.font.size = Pt(18)
    p3.font.bold = True
    p3.font.color.rgb = ACCENT_CYAN
    
    ai_bullets = [
        "Google Gemini API integration",
        "Dynamic prompt engineers",
        "MongoDB Atlas cloud storage",
        "Resume JSON parsing",
        "Auto-portfolio generator"
    ]
    for b in ai_bullets:
        p = tf3.add_paragraph()
        p.text = "🧠 " + b
        p.font.name = "Arial"
        p.font.size = Pt(14)
        p.font.color.rgb = TEXT_LIGHT
        p.space_before = Pt(8)

    # ----------------------------------------------------
    # SLIDE 4: Core Features - Candidate Portal
    # ----------------------------------------------------
    slide4 = prs.slides.add_slide(blank_slide_layout)
    add_slide_header(slide4, "Core Features - Candidate Portal")
    
    # Four feature grids
    feats = [
        ("🧠 AI Career Assessment", "Get customized career suggestions by chatting with an AI Career Advisor that analyzes your interests."),
        ("🗺️ Job & Skill Roadmaps", "Access highly organized step-by-step visual roadmaps for CSE domains like Full Stack, AI, Cyber."),
        ("📄 Resume & Portfolio Builder", "One-click dynamic ATS resume PDF creation and automated live developer portfolio site generation."),
        ("🎤 Mock Interview & Code Analyzer", "Experience AI mock interviews with direct grading and use GitHub analyzer to suggest project improvements.")
    ]
    
    coords = [
        (Inches(0.75), Inches(1.8)),
        (Inches(6.8), Inches(1.8)),
        (Inches(0.75), Inches(4.3)),
        (Inches(6.8), Inches(4.3))
    ]
    
    for i, (title, desc) in enumerate(feats):
        cx, cy = coords[i]
        card = slide4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, cx, cy, Inches(5.7), Inches(2.2))
        card.fill.solid()
        card.fill.fore_color.rgb = ACCENT_BLUE
        card.line.color.rgb = ACCENT_CYAN
        
        tb = slide4.shapes.add_textbox(cx + Inches(0.2), cy + Inches(0.15), Inches(5.3), Inches(1.9))
        tf = tb.text_frame
        tf.word_wrap = True
        
        pt = tf.paragraphs[0]
        pt.text = title
        pt.font.name = "Arial"
        pt.font.size = Pt(18)
        pt.font.bold = True
        pt.font.color.rgb = ACCENT_CYAN
        
        pd = tf.add_paragraph()
        pd.text = desc
        pd.font.name = "Arial"
        pd.font.size = Pt(14)
        pd.font.color.rgb = TEXT_LIGHT
        pd.space_before = Pt(8)

    # ----------------------------------------------------
    # SLIDE 5: Core Features - B2B Recruiter Portal
    # ----------------------------------------------------
    slide5 = prs.slides.add_slide(blank_slide_layout)
    add_slide_header(slide5, "B2B Recruiter Portal")
    
    left_side = slide5.shapes.add_textbox(Inches(0.75), Inches(1.8), Inches(5.6), Inches(4.8))
    tfl = left_side.text_frame
    tfl.word_wrap = True
    
    p = tfl.paragraphs[0]
    p.text = "Recruiter Operations"
    p.font.name = "Arial"
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = ACCENT_CYAN
    
    r_ops = [
        "💡 Smart Candidate Matchmaking: Auto-matches students who have high roadmap completion levels or test scores.",
        "📊 Dynamic HR Analytics Dashboard: Detailed progress statistics, scores, skill sets, and interview grades.",
        "🔍 Single-view Student Dossier: Detailed access to AI-analyzed candidate code quality, resume, and portfolio link.",
        "🎯 Verified Candidate Selection: Minimizes screening times by 80% using pre-screened technical benchmarks."
    ]
    for op in r_ops:
        p = tfl.add_paragraph()
        p.text = op
        p.font.name = "Arial"
        p.font.size = Pt(14)
        p.font.color.rgb = TEXT_LIGHT
        p.space_before = Pt(14)

    # Decorative right card: "Why HR Portal?"
    rcard = slide5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(1.8), Inches(5.7), Inches(4.5))
    rcard.fill.solid()
    rcard.fill.fore_color.rgb = ACCENT_BLUE
    rcard.line.color.rgb = ACCENT_CYAN
    
    tb_rcard = slide5.shapes.add_textbox(Inches(7.1), Inches(2.1), Inches(5.1), Inches(3.9))
    tfr = tb_rcard.text_frame
    tfr.word_wrap = True
    
    pr = tfr.paragraphs[0]
    pr.text = "HR System Core Metrics"
    pr.font.name = "Arial"
    pr.font.size = Pt(20)
    pr.font.bold = True
    pr.font.color.rgb = ACCENT_CYAN
    
    metrics = [
        "⌛ Time-to-Hire Reduction: Cut from 20 days to 3 days",
        "🎯 Profile Matching Accuracy: Up to 92% match index",
        "🌟 Skills-First Hiring: Verified roadmap milestone achievements",
        "🔒 Seamless Integration: HR Dashboard fully synced with candidate database"
    ]
    for m in metrics:
        p = tfr.add_paragraph()
        p.text = m
        p.font.name = "Arial"
        p.font.size = Pt(14)
        p.font.color.rgb = TEXT_LIGHT
        p.space_before = Pt(12)

    # ----------------------------------------------------
    # SLIDE 6: Technical Architecture & Stack
    # ----------------------------------------------------
    slide6 = prs.slides.add_slide(blank_slide_layout)
    add_slide_header(slide6, "Technology Stack & Core API")
    
    # 3 Column stack
    widths = [Inches(3.7), Inches(3.7), Inches(3.7)]
    offsets = [Inches(0.75), Inches(4.8), Inches(8.85)]
    titles = ["FRONTEND TECH", "BACKEND SERVICE", "DATABASE & AI"]
    stacks = [
        ["React 19 Core Framework", "Vite 7 (Ultra-fast bundler)", "React Router 7 (Layout & Routes)", "Vanilla CSS (Neon Glassmorphism)", "TailwindCSS (Responsive utility)"],
        ["Node.js (LTS v18+)", "Express.js Framework", "JWT (Authentication)", "Bcrypt (Password hash protection)", "RESTful API routes"],
        ["MongoDB Atlas (Cloud NoSQL)", "Mongoose (ODM Layer)", "Google Gemini API Integration", "OpenAI (Fallback Service)", "ATS PDF Parser engine"]
    ]
    
    for i in range(3):
        card = slide6.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, offsets[i], Inches(1.8), widths[i], Inches(4.5))
        card.fill.solid()
        card.fill.fore_color.rgb = ACCENT_BLUE
        card.line.color.rgb = ACCENT_CYAN
        
        tb = slide6.shapes.add_textbox(offsets[i] + Inches(0.15), Inches(2.0), widths[i] - Inches(0.3), Inches(4.0))
        tf = tb.text_frame
        tf.word_wrap = True
        
        pt = tf.paragraphs[0]
        pt.text = titles[i]
        pt.font.name = "Arial"
        pt.font.size = Pt(18)
        pt.font.bold = True
        pt.font.color.rgb = ACCENT_CYAN
        
        for item in stacks[i]:
            p = tf.add_paragraph()
            p.text = "✔ " + item
            p.font.name = "Arial"
            p.font.size = Pt(13)
            p.font.color.rgb = TEXT_LIGHT
            p.space_before = Pt(10)

    # ----------------------------------------------------
    # SLIDE 7: Database Design / Schema
    # ----------------------------------------------------
    slide7 = prs.slides.add_slide(blank_slide_layout)
    add_slide_header(slide7, "MongoDB Data Models")
    
    # Left container for Models
    mcard = slide7.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.75), Inches(1.8), Inches(5.6), Inches(4.6))
    mcard.fill.solid()
    mcard.fill.fore_color.rgb = ACCENT_BLUE
    mcard.line.color.rgb = ACCENT_CYAN
    
    tb_m = slide7.shapes.add_textbox(Inches(0.95), Inches(2.0), Inches(5.2), Inches(4.2))
    tf_m = tb_m.text_frame
    tf_m.word_wrap = True
    
    pm = tf_m.paragraphs[0]
    pm.text = "Core Schema Definitions"
    pm.font.name = "Arial"
    pm.font.size = Pt(20)
    pm.font.bold = True
    pm.font.color.rgb = ACCENT_CYAN
    
    schemas = [
        "👤 User Schema: email, password, role ('student'/'recruiter'), profile details, completedRoadmaps, resumes.",
        "📄 Resume Schema: userId, textContent, parsedSkills, scoreATS, generatedPortfolioLink.",
        "🗺️ Roadmap Schema: domainName, level, steps: [{ title, url, desc, completed }].",
        "🎤 MockInterview Schema: userId, topic, score, questionsAnswered: [{ q, a, score, feedback }]."
    ]
    for s in schemas:
        p = tf_m.add_paragraph()
        p.text = s
        p.font.name = "Arial"
        p.font.size = Pt(12)
        p.font.color.rgb = TEXT_LIGHT
        p.space_before = Pt(10)
        
    # Right container for Data Flow
    fcard = slide7.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(1.8), Inches(5.7), Inches(4.6))
    fcard.fill.solid()
    fcard.fill.fore_color.rgb = ACCENT_BLUE
    fcard.line.color.rgb = ACCENT_CYAN
    
    tb_f = slide7.shapes.add_textbox(Inches(7.0), Inches(2.0), Inches(5.3), Inches(4.2))
    tf_f = tb_f.text_frame
    tf_f.word_wrap = True
    
    pf = tf_f.paragraphs[0]
    pf.text = "Data Flow & Integration Flow"
    pf.font.name = "Arial"
    pf.font.size = Pt(20)
    pf.font.bold = True
    pf.font.color.rgb = ACCENT_CYAN
    
    flows = [
        "1. Candidate profiles are completed -> Stored in MongoDB.",
        "2. Resume uploaded -> Sent to Backend -> Parsed using Mongoose & Gemini AI API -> Returns ATS grade + details.",
        "3. Recruiter searches candidates -> Backend queries DB matching specific skills -> Displays matched dossier.",
        "4. Interview completed -> AI assesses speech text -> Grades, saves to DB -> Instantly visible on dashboard."
    ]
    for f in flows:
        p = tf_f.add_paragraph()
        p.text = f
        p.font.name = "Arial"
        p.font.size = Pt(13)
        p.font.color.rgb = TEXT_LIGHT
        p.space_before = Pt(10)

    # ----------------------------------------------------
    # SLIDE 8: Project Highlights & Achievements
    # ----------------------------------------------------
    slide8 = prs.slides.add_slide(blank_slide_layout)
    add_slide_header(slide8, "Project Highlights & Impact")
    
    h_box = slide8.shapes.add_textbox(Inches(0.75), Inches(1.8), Inches(11.83), Inches(4.8))
    tf_h = h_box.text_frame
    tf_h.word_wrap = True
    
    ph1 = tf_h.paragraphs[0]
    ph1.text = "Key System Achievements"
    ph1.font.name = "Arial"
    ph1.font.size = Pt(22)
    ph1.font.bold = True
    ph1.font.color.rgb = ACCENT_CYAN
    
    highlights = [
        "⭐ Real-world Ready: Resolves student career confusion using dynamic AI assessments.",
        "⚡ Automated Workflows: Complete portfolio generation within 2 seconds of profile upload.",
        "📈 Recruiter-Focused: Complete B2B Recruitment module bridging the gap between student portfolios and recruiters.",
        "🔒 Secure & Scalable: Built using best practice patterns like JWT authentication and MongoDB cloud clusters.",
        "🌐 Modern UI Experience: Visually stunning dark mode layout based on modern glassmorphic design principles."
    ]
    for h in highlights:
        p = tf_h.add_paragraph()
        p.text = h
        p.font.name = "Arial"
        p.font.size = Pt(15)
        p.font.color.rgb = TEXT_LIGHT
        p.space_before = Pt(14)

    # ----------------------------------------------------
    # SLIDE 9: Conclusion & Future Scope
    # ----------------------------------------------------
    slide9 = prs.slides.add_slide(blank_slide_layout)
    add_slide_header(slide9, "Conclusion & Future Enhancements")
    
    # Left Box - Conclusion
    box_con = slide9.shapes.add_textbox(Inches(0.75), Inches(1.8), Inches(5.6), Inches(4.5))
    tf_con = box_con.text_frame
    tf_con.word_wrap = True
    p_con = tf_con.paragraphs[0]
    p_con.text = "Conclusion"
    p_con.font.name = "Arial"
    p_con.font.size = Pt(22)
    p_con.font.bold = True
    p_con.font.color.rgb = ACCENT_CYAN
    
    c_points = [
        "CareerCraft bridges the gap between educational preparation and industrial recruitment.",
        "AI Career Assessment provides precise, individualized direction for students.",
        "The B2B Portal provides hiring managers with a highly simplified candidate screening funnel.",
        "Improves student hiring opportunities by showcasing verified milestone accomplishments."
    ]
    for cp in c_points:
        p = tf_con.add_paragraph()
        p.text = "✔ " + cp
        p.font.name = "Arial"
        p.font.size = Pt(15)
        p.font.color.rgb = TEXT_LIGHT
        p.space_before = Pt(12)
        
    # Right Box - Future Scope
    box_fut = slide9.shapes.add_textbox(Inches(6.9), Inches(1.8), Inches(5.6), Inches(4.5))
    tf_fut = box_fut.text_frame
    tf_fut.word_wrap = True
    p_fut = tf_fut.paragraphs[0]
    p_fut.text = "Future Enhancements"
    p_fut.font.name = "Arial"
    p_fut.font.size = Pt(22)
    p_fut.font.bold = True
    p_fut.font.color.rgb = ACCENT_CYAN
    
    f_points = [
        "Integrate live coding environment assessment for candidates.",
        "Add automated voice/video facial expression analysis for mock interviews.",
        "Introduce Multi-Language natural voice guidance chatbots.",
        "Add enterprise blockchain ledger to verify student certificates and project authenticity."
    ]
    for fp in f_points:
        p = tf_fut.add_paragraph()
        p.text = "🚀 " + fp
        p.font.name = "Arial"
        p.font.size = Pt(15)
        p.font.color.rgb = TEXT_LIGHT
        p.space_before = Pt(12)

    # ----------------------------------------------------
    # SLIDE 10: Thank You Slide
    # ----------------------------------------------------
    slide10 = prs.slides.add_slide(blank_slide_layout)
    apply_dark_background(slide10)
    
    # Decorative card
    thank_card = slide10.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.5), Inches(1.5), Inches(10.33), Inches(4.5))
    thank_card.fill.solid()
    thank_card.fill.fore_color.rgb = ACCENT_BLUE
    thank_card.line.color.rgb = ACCENT_CYAN
    thank_card.line.width = Pt(2)
    
    tb_thank = slide10.shapes.add_textbox(Inches(2.0), Inches(2.0), Inches(9.33), Inches(3.5))
    tf_thank = tb_thank.text_frame
    tf_thank.word_wrap = True
    
    pt1 = tf_thank.paragraphs[0]
    pt1.text = "THANK YOU"
    pt1.font.name = "Trebuchet MS"
    pt1.font.size = Pt(54)
    pt1.font.bold = True
    pt1.font.color.rgb = ACCENT_CYAN
    pt1.alignment = PP_ALIGN.CENTER
    
    pt2 = tf_thank.add_paragraph()
    pt2.text = "Questions & Answers Session"
    pt2.font.name = "Arial"
    pt2.font.size = Pt(20)
    pt2.font.color.rgb = TEXT_LIGHT
    pt2.alignment = PP_ALIGN.CENTER
    pt2.space_before = Pt(15)
    
    pt3 = tf_thank.add_paragraph()
    pt3.text = "\nPresented By: HARISHGANTH.L (720823103057)\nHINDUSTAN INSTITUTE OF TECHNOLOGY"
    pt3.font.name = "Arial"
    pt3.font.size = Pt(16)
    pt3.font.color.rgb = TEXT_GRAY
    pt3.alignment = PP_ALIGN.CENTER
    pt3.space_before = Pt(15)
    
    # Save the file
    os.makedirs(os.path.dirname("docs/"), exist_ok=True)
    prs.save("docs/presentation.pptx")
    print("Presentation created successfully at docs/presentation.pptx!")

if __name__ == "__main__":
    create_presentation()
