import os
import pptx
import pymupdf

def inspect_pptx(filepath):
    print(f"--- PPTX Inspection: {filepath} ---")
    if not os.path.exists(filepath):
        print("File does not exist.")
        return
    
    prs = pptx.Presentation(filepath)
    for i, slide in enumerate(prs.slides):
        print(f"\n[Slide {i+1}]")
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                # Check if it has text
                for paragraph in shape.text_frame.paragraphs:
                    text = "".join(run.text for run in paragraph.runs).strip()
                    if text:
                        print(f"  Shape text: {repr(text)}")

def inspect_pdf(filepath):
    print(f"\n--- PDF Inspection: {filepath} ---")
    if not os.path.exists(filepath):
        print("File does not exist.")
        return
    
    doc = pymupdf.open(filepath)
    print(f"Total Pages: {len(doc)}")
    # Print the first few pages content or search for name
    for i in range(min(5, len(doc))):
        print(f"\n[Page {i+1}]")
        text = doc[i].get_text()
        lines = [line.strip() for line in text.split('\n') if line.strip()]
        for line in lines[:20]: # Print first 20 lines
            print(f"  {line}")

if __name__ == "__main__":
    base_dir = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
    pptx_path = os.path.join(base_dir, "docs", "presentation.pptx")
    pdf_path = os.path.join(base_dir, "docs", "report.pdf")
    inspect_pptx(pptx_path)
    inspect_pdf(pdf_path)
